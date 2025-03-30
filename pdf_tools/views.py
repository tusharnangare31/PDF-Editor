from django.shortcuts import render, redirect, get_object_or_404
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.core.files.storage import default_storage
from django.conf import settings
import os
import PyPDF2
from PIL import Image
import pytesseract
import json
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from io import BytesIO
import base64
from django.contrib.auth import login, authenticate, logout
from django.contrib.auth.forms import AuthenticationForm
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.urls import reverse
from .models import ProcessedDocument, DocumentHistory, UserProfile
from .forms import CustomUserCreationForm
import time
import tempfile
import subprocess
from pathlib import Path
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.lib.pagesizes import letter
import uuid

def home(request):
    return render(request, 'pdf_tools/home.html')

def organize(request):
    return render(request, 'pdf_tools/organize.html')

def convert(request):
    return render(request, 'pdf_tools/convert.html')

def edit(request):
    return render(request, 'pdf_tools/edit.html')

def security(request):
    return render(request, 'pdf_tools/security.html')

# Individual feature views
def merge_pdf_view(request):
    return render(request, 'pdf_tools/features/merge.html')

def split_pdf_view(request):
    return render(request, 'pdf_tools/features/split.html')

def compress_pdf_view(request):
    return render(request, 'pdf_tools/features/compress.html')

def jpg_to_pdf_view(request):
    return render(request, 'pdf_tools/features/jpg-to-pdf.html')

# New view functions for feature pages
def ocr_pdf_view(request):
    return render(request, 'pdf_tools/features/ocr.html')

def pdf_to_jpg_view(request):
    return render(request, 'pdf_tools/features/pdf-to-jpg.html')

def pdf_to_word_view(request):
    return render(request, 'pdf_tools/features/pdf-to-word.html')

def pdf_to_ppt_view(request):
    return render(request, 'pdf_tools/features/pdf-to-ppt.html')

def word_to_pdf_view(request):
    return render(request, 'pdf_tools/features/word-to-pdf.html')

def ppt_to_pdf_view(request):
    return render(request, 'pdf_tools/features/ppt-to-pdf.html')

def rotate_pdf_view(request):
    return render(request, 'pdf_tools/features/rotate.html')

def add_page_numbers_view(request):
    return render(request, 'pdf_tools/features/add-page-numbers.html')

def add_watermark_view(request):
    return render(request, 'pdf_tools/features/add-watermark.html')

def protect_pdf_view(request):
    return render(request, 'pdf_tools/features/protect.html')

def unlock_pdf_view(request):
    return render(request, 'pdf_tools/features/unlock.html')

def remove_pages_view(request):
    return render(request, 'pdf_tools/features/remove-pages.html')

@csrf_exempt
def merge_pdf(request):
    if request.method == 'POST':
        try:
            files = request.FILES.getlist('files')
            if not files:
                return JsonResponse({
                    'success': False,
                    'error': 'No files provided'
                })
                
            merger = PyPDF2.PdfMerger()
            
            for file in files:
                merger.append(file)
            
            output_path = os.path.join(settings.MEDIA_ROOT, 'merged.pdf')
            with open(output_path, 'wb') as output_file:
                merger.write(output_file)
            
            # Save to user's documents if authenticated
            document = None
            if request.user.is_authenticated:
                file_names = [f.name for f in files]
                file_info = ", ".join(file_names[:3])
                if len(files) > 3:
                    file_info += f" and {len(files) - 3} more"
                
                document = save_processed_file(
                    user=request.user,
                    file_path=output_path,
                    title=f"Merged PDF ({len(files)} files)",
                    document_type='merged',
                    operation='merge',
                    details=json.dumps({
                        'files': file_info,
                        'count': len(files)
                    })
                )
            
            file_url = f'/media/merged.pdf'
            if document:
                file_url = document.file.url
                
            return JsonResponse({
                'success': True,
                'file_url': file_url
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def split_pdf(request):
    if request.method == 'POST':
        try:
            print("Split PDF request received")
            
            # Check if file exists in request
            if 'file' not in request.FILES:
                print("Error: No file in request")
                return JsonResponse({
                    'success': False,
                    'error': 'No file provided. Please select a PDF file to split.'
                })
            
            file = request.FILES['file']
            print(f"File received: {file.name}, size: {file.size} bytes")
            
            # Validate file is a PDF
            if not file.name.lower().endswith('.pdf'):
                print("Error: File is not a PDF")
                return JsonResponse({
                    'success': False,
                    'error': 'The uploaded file is not a PDF. Please select a valid PDF file.'
                })
            
            # Try to read the PDF header
            try:
                pdf_header = file.read(5)
                file.seek(0)  # Reset file pointer
                if not pdf_header.startswith(b'%PDF-'):
                    print("Error: File does not have PDF header")
                    return JsonResponse({
                        'success': False,
                        'error': 'The uploaded file is not a valid PDF. It does not have a proper PDF header.'
                    })
            except Exception as e:
                print(f"Error checking PDF header: {str(e)}")
                # Continue anyway, let PyPDF2 validate the PDF
            
            # Check if pages parameter exists
            if 'pages' not in request.POST:
                print("Error: No pages parameter")
                return JsonResponse({
                    'success': False,
                    'error': 'No pages specified. Please select which pages to extract.'
                })
            
            # Get the pages input string
            pages_input = request.POST.get('pages', '')
            print(f"Raw pages input: {pages_input}")
            
            # Create a temporary file to handle the uploaded PDF
            temp_dir = tempfile.mkdtemp()
            input_path = os.path.join(temp_dir, 'input.pdf')
            
            # Save the uploaded file
            with open(input_path, 'wb+') as destination:
                for chunk in file.chunks():
                    destination.write(chunk)
            
            # Try to parse the pages
            clean_pages = []
            
            try:
                # First, try to parse as JSON if it looks like a JSON array
                if pages_input.strip().startswith('[') and pages_input.strip().endswith(']'):
                    try:
                        pages_list = json.loads(pages_input)
                        if isinstance(pages_list, list):
                            # Convert all elements to integers
                            clean_pages = [int(p) for p in pages_list if str(p).isdigit()]
                            print(f"Parsed JSON array to: {clean_pages}")
                    except json.JSONDecodeError:
                        print("Input looks like JSON but couldn't be parsed as JSON, falling back to string parsing")
                
                # If we don't have pages yet, try string parsing
                if not clean_pages:
                    # Load the PDF to get total pages
                    try:
                        reader = PyPDF2.PdfReader(input_path)
                        total_pages = len(reader.pages)
                        print(f"PDF has {total_pages} pages")
                        
                        # Handle special keywords
                        if pages_input.strip() == 'all':
                            clean_pages = list(range(1, total_pages + 1))
                        elif pages_input.strip() == 'allbutfirst' and total_pages > 1:
                            clean_pages = list(range(2, total_pages + 1))
                        elif pages_input.strip() == 'allbutlast' and total_pages > 1:
                            clean_pages = list(range(1, total_pages))
                        elif pages_input.strip() == 'oddpages':
                            clean_pages = list(range(1, total_pages + 1, 2))
                        elif pages_input.strip() == 'evenpages':
                            clean_pages = list(range(2, total_pages + 1, 2))
                        else:
                            # Parse comma-separated values and ranges
                            parts = pages_input.split(',')
                            
                            for part in parts:
                                part = part.strip()
                                if not part:
                                    continue
                                    
                                # Handle ranges like "1-5"
                                if '-' in part:
                                    range_parts = part.split('-')
                                    
                                    # Check if this is an open-ended range like "2-"
                                    if len(range_parts) == 2 and range_parts[1].strip() == '':
                                        # Open-ended range starting at some number
                                        try:
                                            start = int(range_parts[0].strip())
                                            if start >= 1 and start <= total_pages:
                                                clean_pages.extend(range(start, total_pages + 1))
                                        except ValueError:
                                            print(f"Invalid open-ended range start: {range_parts[0]}")
                                        continue
                                    
                                    # Normal range with start and end
                                    if len(range_parts) == 2:
                                        try:
                                            start = int(range_parts[0].strip())
                                            end = int(range_parts[1].strip())
                                            
                                            if start >= 1 and end <= total_pages and start <= end:
                                                clean_pages.extend(range(start, end + 1))
                                            else:
                                                print(f"Range out of bounds or invalid: {start}-{end}, PDF has {total_pages} pages")
                                        except ValueError:
                                            print(f"Invalid range: {part}")
                                else:
                                    # Single page number
                                    try:
                                        page_num = int(part)
                                        if page_num >= 1 and page_num <= total_pages:
                                            clean_pages.append(page_num)
                                        else:
                                            print(f"Page number out of bounds: {page_num}, PDF has {total_pages} pages")
                                    except ValueError:
                                        print(f"Invalid page number: {part}")
                    except Exception as pdf_error:
                        print(f"Error reading PDF: {str(pdf_error)}")
                        return JsonResponse({
                            'success': False,
                            'error': f'Error reading the PDF file: {str(pdf_error)}'
                        })
                        
                    # Remove duplicates and sort
                    clean_pages = sorted(list(set(clean_pages)))
                    print(f"Parsed page string to: {clean_pages}")
            except Exception as parse_error:
                print(f"Error parsing pages: {str(parse_error)}")
                import traceback
                traceback.print_exc()
                return JsonResponse({
                    'success': False,
                    'error': f'Error parsing page numbers: {str(parse_error)}'
                })
            
            # Verify we have valid pages after cleaning
            if not clean_pages:
                print("No valid pages after parsing")
                return JsonResponse({
                    'success': False,
                    'error': 'No valid pages found. Please check your page range and ensure it is within the document limits.'
                })
            
            print(f"Final pages to extract: {clean_pages}")
            
            try:
                # Create PDF reader and writer
                reader = PyPDF2.PdfReader(input_path)
                writer = PyPDF2.PdfWriter()
                
                # Get total pages for validation
                total_pages = len(reader.pages)
                print(f"PDF has {total_pages} pages")
                
                # Store the original requested pages before filtering
                requested_pages = clean_pages.copy()
                
                # Filter out any pages that are out of bounds
                valid_pages = [p for p in clean_pages if p >= 1 and p <= total_pages]
                skipped_pages = [p for p in clean_pages if p < 1 or p > total_pages]
                
                if len(valid_pages) < len(clean_pages):
                    print(f"Filtered out {len(clean_pages) - len(valid_pages)} pages that were out of bounds")
                
                # Verify we still have pages to extract
                if not valid_pages:
                    return JsonResponse({
                        'success': False,
                        'error': f'No valid pages within document range. The PDF only has {total_pages} pages, but you requested: {", ".join(map(str, requested_pages))}'
                    })
                
                # Add pages to the writer
                processed_pages = []
                failed_pages = []
                for page_num in valid_pages:
                    try:
                        print(f"Adding page {page_num}")
                        # PyPDF2 uses 0-based indexing
                        page = reader.pages[page_num - 1]
                        writer.add_page(page)
                        processed_pages.append(page_num)
                    except Exception as page_error:
                        print(f"Error adding page {page_num}: {str(page_error)}")
                        import traceback
                        traceback.print_exc()
                        # Track failed pages
                        failed_pages.append(page_num)
                        # Continue with other pages instead of failing completely
                        continue
                
                # If no pages were successfully processed
                if not processed_pages:
                    return JsonResponse({
                        'success': False,
                        'error': 'Failed to process any of the specified pages. Please check your PDF file integrity.'
                    })
                
                # Ensure media directory exists
                os.makedirs(settings.MEDIA_ROOT, exist_ok=True)
                
                # Create a unique filename to avoid overwrites
                timestamp = int(time.time())
                output_filename = f'split_{timestamp}.pdf'
                output_path = os.path.join(settings.MEDIA_ROOT, output_filename)
                
                # Make sure parent directory exists
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                print(f"Writing output to: {output_path}")
                
                # Write the output file
                with open(output_path, 'wb') as output_file:
                    writer.write(output_file)
                
                print("PDF successfully split")
                
                # Clean up temp directory
                try:
                    import shutil
                    shutil.rmtree(temp_dir)
                except Exception as e:
                    print(f"Error cleaning up temp directory: {str(e)}")
                
                # Save to user's documents if authenticated
                document = None
                if request.user.is_authenticated:
                    print(f"User is authenticated: {request.user.username}")
                    page_info = f"Pages {', '.join(map(str, processed_pages))}"
                    
                    try:
                        document = save_processed_file(
                            user=request.user,
                            file_path=output_path,
                            title=f"Split PDF ({len(processed_pages)} pages)",
                            document_type='split',
                            operation='split',
                            details=json.dumps({
                                'pages': page_info,
                                'count': len(processed_pages),
                                'original_filename': file.name
                            })
                        )
                        print(f"Document saved to user's account: {document.id if document else 'None'}")
                    except Exception as save_error:
                        print(f"Error saving document to user account: {str(save_error)}")
                        # Continue even if saving to user account fails
                else:
                    print("User is not authenticated, not saving to account")
                
                file_url = f'/media/{output_filename}'
                if document:
                    file_url = document.file.url
                
                # Create a status message for the frontend
                status_message = ""
                if skipped_pages:
                    status_message = f"Note: {len(skipped_pages)} page(s) were skipped because they exceed the document's page count ({total_pages} pages): {', '.join(map(str, skipped_pages))}"
                
                print(f"Returning success with file URL: {file_url}")
                return JsonResponse({
                    'success': True,
                    'file_url': file_url,
                    'pages_extracted': processed_pages,
                    'total_pdf_pages': total_pages,
                    'requested_pages': requested_pages,
                    'skipped_pages': skipped_pages,
                    'failed_pages': failed_pages,
                    'total_pages': len(processed_pages),
                    'status_message': status_message
                })
            except Exception as pdf_error:
                print(f"PDF processing error: {str(pdf_error)}")
                import traceback
                traceback.print_exc()
                return JsonResponse({
                    'success': False,
                    'error': f'PDF processing error: {str(pdf_error)}'
                })
                
        except Exception as e:
            print(f"Unexpected error in split_pdf: {str(e)}")
            import traceback
            traceback.print_exc()
            return JsonResponse({
                'success': False,
                'error': f'Unexpected error: {str(e)}'
            })
    
    print("Invalid request method for split_pdf")
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def compress_pdf(request):
    if request.method == 'POST':
        try:
            print("Compression request received")
            
            # Check if file exists in request
            if 'file' not in request.FILES:
                print("Error: No file in request")
                return JsonResponse({
                    'success': False,
                    'error': 'No file was uploaded. Please select a file.'
                })
                
            file = request.FILES['file']
            compression_level = request.POST.get('compressionLevel', 'medium')  # Default to medium compression
            
            print(f"File received: {file.name}, size: {file.size} bytes, compression level: {compression_level}")
            
            # Check if the file is a valid PDF
            if not file.name.lower().endswith('.pdf'):
                print("Error: Not a PDF file")
                return JsonResponse({
                    'success': False,
                    'error': 'The uploaded file is not a PDF. Please select a valid PDF file.'
                })
            
            # Check first few bytes for PDF signature
            try:
                pdf_header = file.read(5)
                file.seek(0)  # Reset file pointer
                if not pdf_header.startswith(b'%PDF-'):
                    print("Error: File does not have PDF header")
                    return JsonResponse({
                        'success': False,
                        'error': 'The uploaded file is not a valid PDF. It does not have a proper PDF header.'
                    })
            except Exception as e:
                print(f"Error checking PDF header: {str(e)}")
                # Continue anyway, let PyPDF2 validate the PDF
            
            # Create a temp directory to work with files
            temp_dir = tempfile.mkdtemp()
            input_path = os.path.join(temp_dir, 'input.pdf')
            output_path = os.path.join(temp_dir, 'output.pdf')
            timestamp = int(time.time())
            final_path = os.path.join(settings.MEDIA_ROOT, f'compressed_{timestamp}.pdf')
            
            # Ensure media directory exists
            os.makedirs(settings.MEDIA_ROOT, exist_ok=True)
            
            # Save the uploaded file to temp directory
            with open(input_path, 'wb') as f:
                for chunk in file.chunks():
                    f.write(chunk)
            
            # Verify the file is a valid PDF using PyPDF2
            try:
                reader = PyPDF2.PdfReader(input_path)
                page_count = len(reader.pages)
                print(f"Valid PDF with {page_count} pages")
                if page_count == 0:
                    return JsonResponse({
                        'success': False,
                        'error': 'The uploaded PDF has no pages to compress.'
                    })
            except Exception as e:
                print(f"Error validating PDF: {str(e)}")
                return JsonResponse({
                    'success': False,
                    'error': f'The uploaded file is not a valid PDF. Error: {str(e)}'
                })
            
            # Get original file size
            original_size = os.path.getsize(input_path)
            print(f"Original file size: {original_size} bytes")
            
            compressed_size = original_size
            compression_success = False
            compression_method = "original"
            
            # Try to use pikepdf for better compression
            try:
                print("Attempting to use pikepdf for compression")
                import pikepdf
                
                # Compression settings based on level
                if compression_level == 'low':
                    # Low compression - preserve quality
                    image_dpi = 150
                    jpeg_quality = 85
                    compress_streams = True
                elif compression_level == 'high':
                    # High compression - lower quality
                    image_dpi = 72
                    jpeg_quality = 30
                    compress_streams = True
                else:  # medium (default)
                    # Medium compression - balanced
                    image_dpi = 100
                    jpeg_quality = 60
                    compress_streams = True
                
                # Use pikepdf to compress
                pdf = pikepdf.open(input_path)
                
                print(f"Using compression settings: DPI={image_dpi}, JPEG quality={jpeg_quality}")
                
                # Save with compression settings
                pdf.save(output_path, 
                    compress_streams=compress_streams,
                    object_stream_mode=pikepdf.ObjectStreamMode.generate,
                    compress_content_streams=True,
                    recompress_flate=True,
                    recompress_jpeg=True,
                    jpeg_quality=jpeg_quality)
                
                pdf.close()
                
                if os.path.exists(output_path):
                    compressed_size = os.path.getsize(output_path)
                    print(f"Pikepdf compressed size: {compressed_size} bytes")
                    compression_success = True
                    compression_method = "pikepdf"
                else:
                    print("Pikepdf output file not found")
                    raise Exception("Failed to create compressed PDF with pikepdf")
                
                # If compression was not effective, try more aggressive optimization with Ghostscript-like settings
                if compressed_size > original_size * 0.95:
                    print("Initial compression not very effective, trying more aggressive settings")
                    
                    # More aggressive settings
                    pdf = pikepdf.open(input_path)
                    
                    # Try to optimize images more aggressively
                    from pikepdf import PdfImage
                    
                    for i, page in enumerate(pdf.pages):
                        for name, obj in page.images.items():
                            try:
                                img = PdfImage(obj)
                                if compression_level == 'high':
                                    # For high compression, attempt to extract, compress, and replace images
                                    if img.bits_per_component > 4:
                                        try:
                                            print(f"Processing image {name} on page {i+1}")
                                            # Skip if image is already small
                                            if obj.get_raw_stream_buffer().size < 5000:
                                                continue
                                            
                                            # This advanced image replacement would normally go here
                                            # For now, we'll rely on pikepdf's built-in compression
                                        except Exception as e:
                                            print(f"Error processing image: {str(e)}")
                            except Exception as e:
                                print(f"Error extracting image: {str(e)}")
                    
                    # Save with more aggressive settings
                    pdf.save(output_path, 
                        compress_streams=True,
                        object_stream_mode=pikepdf.ObjectStreamMode.generate,
                        stream_decode_level=pikepdf.StreamDecodeLevel.generalized,
                        compress_content_streams=True,
                        recompress_flate=True,
                        recompress_jpeg=True,
                        jpeg_quality=max(20, jpeg_quality-20))  # Even lower quality
                    
                    pdf.close()
                    
                    if os.path.exists(output_path):
                        new_compressed_size = os.path.getsize(output_path)
                        print(f"Aggressive pikepdf compressed size: {new_compressed_size} bytes")
                        
                        # Only use the aggressive compression if it's actually better
                        if new_compressed_size < compressed_size:
                            compressed_size = new_compressed_size
                            compression_method = "pikepdf-aggressive"
                    else:
                        print("Aggressive compression output file not found")
                
                # Copy the result to the final destination
                import shutil
                shutil.copy(output_path, final_path)
            
            except ImportError:
                print("pikepdf not available, falling back to PyPDF2")
                # Fall back to PyPDF2 if pikepdf is not available
                compression_success = False
                
                try:
                    # Try to use PyPDF2 compression
                    reader = PyPDF2.PdfReader(input_path)
                    writer = PyPDF2.PdfWriter()
                    
                    for page in reader.pages:
                        try:
                            page.compress_content_streams()  # Basic compression
                        except Exception as e:
                            print(f"Error compressing page: {str(e)}")
                        writer.add_page(page)
                    
                    # PyPDF2 compression
                    writer.compress = True
                    
                    with open(output_path, 'wb') as output_file:
                        writer.write(output_file)
                    
                    if os.path.exists(output_path):
                        pdfwrite_size = os.path.getsize(output_path)
                        compressed_size = pdfwrite_size
                        
                        if compressed_size < original_size:
                            compression_success = True
                            compression_method = "pypdf2"
                        
                        # Copy if smaller
                        import shutil
                        shutil.copy(output_path, final_path)
                except Exception as e:
                    print(f"Error with PyPDF2 compression: {str(e)}")
                    # If PyPDF2 compression fails, use the original file
                    import shutil
                    shutil.copy(input_path, final_path)
            
            except Exception as e:
                print(f"Error with pikepdf compression: {str(e)}")
                import traceback
                traceback.print_exc()
                # Fall back to PyPDF2
                try:
                    # Try to use PyPDF2 compression
                    reader = PyPDF2.PdfReader(input_path)
                    writer = PyPDF2.PdfWriter()
                    
                    for page in reader.pages:
                        try:
                            page.compress_content_streams()  # Basic compression
                        except Exception as e:
                            print(f"Error compressing page: {str(e)}")
                        writer.add_page(page)
                    
                    # PyPDF2 compression
                    writer.compress = True
                    
                    with open(output_path, 'wb') as output_file:
                        writer.write(output_file)
                    
                    if os.path.exists(output_path):
                        pdfwrite_size = os.path.getsize(output_path)
                        compressed_size = pdfwrite_size
                        
                        if compressed_size < original_size:
                            compression_success = True
                            compression_method = "pypdf2"
                        
                        # Copy if smaller
                        import shutil
                        shutil.copy(output_path, final_path)
                    else:
                        # If compression fails, use the original file
                        import shutil
                        compressed_size = original_size
                        shutil.copy(input_path, final_path)
                except Exception as e2:
                    print(f"Error with PyPDF2 compression: {str(e2)}")
                    # If PyPDF2 compression fails, use the original file
                    import shutil
                    compressed_size = original_size
                    shutil.copy(input_path, final_path)
            
            # Make sure we have a final compressed file
            if not os.path.exists(final_path):
                print("Final compressed file does not exist, copying original")
                import shutil
                shutil.copy(input_path, final_path)
                compressed_size = original_size
            else:
                compressed_size = os.path.getsize(final_path)
            
            # If compression made file larger, use original
            if compressed_size >= original_size:
                print(f"Compression increased file size ({compressed_size} > {original_size}), using original")
                import shutil
                shutil.copy(input_path, final_path)
                compressed_size = original_size
                size_reduction = 0
                compression_method = "original"
            else:
                # Calculate size reduction percentage
                size_reduction = round(((original_size - compressed_size) / original_size) * 100)
                print(f"Compression successful: {original_size} -> {compressed_size} bytes ({size_reduction}% reduction) using {compression_method}")
            
            # Save to user's documents if authenticated
            document = None
            if request.user.is_authenticated:
                document = save_processed_file(
                    user=request.user,
                    file_path=final_path,
                    title=f"Compressed PDF ({compression_level} compression)",
                    document_type='compressed',
                    operation='compress',
                    details=json.dumps({
                        'original_size': original_size,
                        'compressed_size': compressed_size,
                        'size_reduction': size_reduction,
                        'compression_level': compression_level,
                        'compression_method': compression_method,
                        'original_filename': file.name
                    })
                )
            
            # Clean up temp files
            try:
                import shutil
                shutil.rmtree(temp_dir)
            except Exception as e:
                print(f"Error cleaning up temp files: {str(e)}")
            
            # Construct the URL for the compressed file
            file_url = f'/media/{os.path.basename(final_path)}'
            if document:
                file_url = document.file.url
            
            return JsonResponse({
                'success': True,
                'file_url': file_url,
                'original_size': original_size,
                'compressed_size': compressed_size,
                'size_reduction': size_reduction,
                'compression_method': compression_method
            })
        except Exception as e:
            print(f"Error in compress_pdf: {str(e)}")
            import traceback
            traceback.print_exc()
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def ocr_pdf(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            reader = PyPDF2.PdfReader(file)
            writer = PyPDF2.PdfWriter()
            
            for page in reader.pages:
                # Convert PDF page to image
                image = page.convert_to_image()
                # Perform OCR
                text = pytesseract.image_to_string(image)
                # Add text as annotation
                page.add_annotation(text)
                writer.add_page(page)
            
            output_path = os.path.join(settings.MEDIA_ROOT, 'ocr.pdf')
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/ocr.pdf'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def jpg_to_pdf(request):
    if request.method == 'POST':
        try:
            files = request.FILES.getlist('files')
            page_layout = request.POST.get('pageLayout', '1')  # Default to 1 image per page
            
            # Convert page_layout to int
            images_per_page = int(page_layout)
            
            # Create a PDF with multiple images
            output_path = os.path.join(settings.MEDIA_ROOT, 'converted.pdf')
            
            if len(files) == 1 and images_per_page == 1:
                # Simple case - just one image
                image = Image.open(files[0])
                
                # Convert to RGB if necessary
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                image.save(output_path, 'PDF')
            else:
                # Multiple images or specific layout
                first_image = Image.open(files[0])
                if first_image.mode != 'RGB':
                    first_image = first_image.convert('RGB')
                
                # Calculate page size based on first image
                width, height = first_image.size
                
                # Create a list to hold all images
                images = [first_image]
                
                # Open and convert all other images
                for i in range(1, len(files)):
                    img = Image.open(files[i])
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    images.append(img)
                
                # Create new images for multi-image pages if needed
                if images_per_page > 1:
                    final_images = []
                    for i in range(0, len(images), images_per_page):
                        batch = images[i:i+images_per_page]
                        if len(batch) == images_per_page or images_per_page == 4:
                            # Create a new image to hold the batch
                            if images_per_page == 2:
                                # 2 images per page - side by side
                                new_img = Image.new('RGB', (width * 2, height), (255, 255, 255))
                                new_img.paste(batch[0], (0, 0))
                                if len(batch) > 1:
                                    new_img.paste(batch[1], (width, 0))
                            else:  # 4 images per page
                                # 4 images per page - 2x2 grid
                                new_img = Image.new('RGB', (width * 2, height * 2), (255, 255, 255))
                                new_img.paste(batch[0], (0, 0))
                                if len(batch) > 1:
                                    new_img.paste(batch[1], (width, 0))
                                if len(batch) > 2:
                                    new_img.paste(batch[2], (0, height))
                                if len(batch) > 3:
                                    new_img.paste(batch[3], (width, height))
                            final_images.append(new_img)
                        else:
                            # Just add the remaining images individually
                            final_images.extend(batch)
                    
                    # Save the final images
                    if final_images:
                        final_images[0].save(output_path, 'PDF', save_all=True, append_images=final_images[1:])
                else:
                    # Save all images to PDF (1 image per page)
                    images[0].save(output_path, 'PDF', save_all=True, append_images=images[1:])
            
            # Save to user's documents if authenticated
            document = None
            if request.user.is_authenticated:
                file_names = [f.name for f in files]
                file_info = ", ".join(file_names[:3])
                if len(files) > 3:
                    file_info += f" and {len(files) - 3} more"
                
                document = save_processed_file(
                    user=request.user,
                    file_path=output_path,
                    title=f"JPG to PDF ({len(files)} images)",
                    document_type='converted',
                    operation='convert_to',
                    details=json.dumps({
                        'files': file_info,
                        'count': len(files),
                        'layout': f"{images_per_page} per page"
                    })
                )
            
            file_url = f'/media/converted.pdf'
            if document:
                file_url = document.file.url
            
            return JsonResponse({
                'success': True,
                'file_url': file_url
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def word_to_pdf(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            doc = Document(file)
            
            # Convert to PDF using a PDF library
            # This is a simplified example - you might want to use a more robust solution
            output_path = os.path.join(settings.MEDIA_ROOT, 'converted.pdf')
            # Add your Word to PDF conversion logic here
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/converted.pdf'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def ppt_to_pdf(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            prs = Presentation(file)
            
            # Convert to PDF using a PDF library
            # This is a simplified example - you might want to use a more robust solution
            output_path = os.path.join(settings.MEDIA_ROOT, 'converted.pdf')
            # Add your PowerPoint to PDF conversion logic here
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/converted.pdf'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def pdf_to_jpg(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            images = convert_from_path(file)
            
            # Save the first page as JPG
            output_path = os.path.join(settings.MEDIA_ROOT, 'converted.jpg')
            images[0].save(output_path, 'JPEG')
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/converted.jpg'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def pdf_to_word(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            reader = PyPDF2.PdfReader(file)
            
            # Extract text from PDF
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            
            # Create Word document
            doc = Document()
            doc.add_paragraph(text)
            
            output_path = os.path.join(settings.MEDIA_ROOT, 'converted.docx')
            doc.save(output_path)
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/converted.docx'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def pdf_to_ppt(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            reader = PyPDF2.PdfReader(file)
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Add each page as a slide
            for page in reader.pages:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                # Add content to slide
                # This is a simplified example - you might want to add more formatting
            
            output_path = os.path.join(settings.MEDIA_ROOT, 'converted.pptx')
            prs.save(output_path)
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/converted.pptx'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def rotate_pdf(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            angle = int(request.POST['angle'])
            pages = request.POST['pages']
            
            reader = PyPDF2.PdfReader(file)
            writer = PyPDF2.PdfWriter()
            
            # Parse page range
            page_numbers = []
            for part in pages.split(','):
                if '-' in part:
                    start, end = map(int, part.split('-'))
                    page_numbers.extend(range(start, end + 1))
                else:
                    page_numbers.append(int(part))
            
            # Rotate specified pages
            for i, page in enumerate(reader.pages):
                if i + 1 in page_numbers:
                    page.rotate(angle)
                writer.add_page(page)
            
            output_path = os.path.join(settings.MEDIA_ROOT, 'rotated.pdf')
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/rotated.pdf'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def add_page_numbers(request):
    if request.method == 'POST':
        try:
            import tempfile
            import os
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import inch
            from reportlab.lib.pagesizes import letter
            from io import BytesIO
            from PyPDF2 import PdfReader, PdfWriter
            import uuid
            
            file = request.FILES['file']
            position = request.POST['position']
            start_number = int(request.POST['start_number'])
            format_type = request.POST.get('format', 'numeric')
            skip_first_page = request.POST.get('skip_first_page') == 'true'
            
            # Create a temporary directory
            temp_dir = tempfile.mkdtemp()
            
            # Save uploaded file to temp directory
            input_path = os.path.join(temp_dir, f"input_{uuid.uuid4()}.pdf")
            with open(input_path, 'wb+') as destination:
                for chunk in file.chunks():
                    destination.write(chunk)
            
            # Read the PDF
            reader = PdfReader(input_path)
            writer = PdfWriter()
            total_pages = len(reader.pages)
            
            # Position mapping
            positions = {
                'bottom-center': {'x': 0.5, 'y': 0.5, 'align': 'center'},
                'bottom-right': {'x': 0.9, 'y': 0.5, 'align': 'right'},
                'bottom-left': {'x': 0.1, 'y': 0.5, 'align': 'left'},
                'top-center': {'x': 0.5, 'y': 10.5, 'align': 'center'},
                'top-right': {'x': 0.9, 'y': 10.5, 'align': 'right'},
                'top-left': {'x': 0.1, 'y': 10.5, 'align': 'left'}
            }
            
            # Get position coordinates
            pos = positions.get(position, positions['bottom-center'])
            
            # Create page numbers for each page
            for i in range(total_pages):
                # Skip first page if requested
                if skip_first_page and i == 0:
                    writer.add_page(reader.pages[i])
                    continue
                
                # Calculate page number
                page_number = start_number + i
                
                # Format page number based on selected format
                if format_type == 'roman':
                    # Convert to roman numerals
                    def to_roman(num):
                        val = [
                            1000, 900, 500, 400,
                            100, 90, 50, 40,
                            10, 9, 5, 4,
                            1
                        ]
                        syms = [
                            "M", "CM", "D", "CD",
                            "C", "XC", "L", "XL",
                            "X", "IX", "V", "IV",
                            "I"
                        ]
                        roman_num = ''
                        i = 0
                        while num > 0:
                            for _ in range(num // val[i]):
                                roman_num += syms[i]
                                num -= val[i]
                            i += 1
                        return roman_num
                    
                    page_text = to_roman(page_number)
                elif format_type == 'page-of-total':
                    page_text = f"Page {page_number} of {total_pages}"
                else:  # numeric
                    page_text = str(page_number)
                
                # Get original page
                original_page = reader.pages[i]
                page_width = float(original_page.mediabox.width)
                page_height = float(original_page.mediabox.height)
                
                # Create a PDF with just the page number
                packet = BytesIO()
                c = canvas.Canvas(packet, pagesize=(page_width, page_height))
                
                # Set font and size
                c.setFont("Helvetica", 10)
                
                # Calculate text position
                if pos['align'] == 'center':
                    x_pos = page_width * pos['x']
                    c.drawCentredString(x_pos, page_height / 12 if 'bottom' in position else page_height - 20, page_text)
                elif pos['align'] == 'right':
                    x_pos = page_width * pos['x']
                    c.drawRightString(x_pos, page_height / 12 if 'bottom' in position else page_height - 20, page_text)
                else:  # left alignment
                    x_pos = page_width * pos['x']
                    c.drawString(x_pos, page_height / 12 if 'bottom' in position else page_height - 20, page_text)
                
                c.save()
                
                # Move to the beginning of the StringIO buffer
                packet.seek(0)
                
                # Create a new PDF with the page number
                number_pdf = PdfReader(packet)
                
                # Merge original page with the page number
                page = reader.pages[i]
                page.merge_page(number_pdf.pages[0])
                
                # Add the page with numbers to the writer
                writer.add_page(page)
            
            # Save the numbered PDF
            output_filename = f"numbered_{uuid.uuid4()}.pdf"
            output_path = os.path.join(settings.MEDIA_ROOT, output_filename)
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            # Clean up temporary files
            try:
                os.remove(input_path)
                os.rmdir(temp_dir)
            except:
                pass
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/{output_filename}'
            })
        except Exception as e:
            import traceback
            print(traceback.format_exc())
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def add_watermark(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            text = request.POST['text']
            opacity = float(request.POST['opacity']) / 100
            
            reader = PyPDF2.PdfReader(file)
            writer = PyPDF2.PdfWriter()
            
            for page in reader.pages:
                # Add watermark text
                # This is a simplified example - you might want to add more formatting
                writer.add_page(page)
            
            output_path = os.path.join(settings.MEDIA_ROOT, 'watermarked.pdf')
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/watermarked.pdf'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def unlock_pdf(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            password = request.POST['password']
            
            reader = PyPDF2.PdfReader(file)
            writer = PyPDF2.PdfWriter()
            
            # Try to decrypt with password
            if reader.is_encrypted:
                reader.decrypt(password)
            
            for page in reader.pages:
                writer.add_page(page)
            
            output_path = os.path.join(settings.MEDIA_ROOT, 'unlocked.pdf')
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/unlocked.pdf'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def protect_pdf(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            password = request.POST['password']
            restrict_printing = request.POST.get('restrict_printing') == 'true'
            restrict_copying = request.POST.get('restrict_copying') == 'true'
            restrict_editing = request.POST.get('restrict_editing') == 'true'
            
            reader = PyPDF2.PdfReader(file)
            writer = PyPDF2.PdfWriter()
            
            for page in reader.pages:
                writer.add_page(page)
            
            # Set encryption
            writer.encrypt(password)
            
            output_path = os.path.join(settings.MEDIA_ROOT, 'protected.pdf')
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/protected.pdf'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def sign_pdf(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            signature = request.POST['signature']
            
            # Decode base64 signature
            signature_data = base64.b64decode(signature.split(',')[1])
            signature_image = Image.open(BytesIO(signature_data))
            
            reader = PyPDF2.PdfReader(file)
            writer = PyPDF2.PdfWriter()
            
            for page in reader.pages:
                # Add signature to the last page
                if page == reader.pages[-1]:
                    # Add signature image
                    # This is a simplified example - you might want to add more positioning
                    writer.add_page(page)
            
            output_path = os.path.join(settings.MEDIA_ROOT, 'signed.pdf')
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/signed.pdf'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

@csrf_exempt
def remove_pages(request):
    if request.method == 'POST':
        try:
            file = request.FILES['file']
            # The frontend sends 'page_range' but we were looking for 'pages'
            pages_to_remove = request.POST.get('page_range')
            
            if not pages_to_remove:
                return JsonResponse({
                    'success': False,
                    'error': 'Page range is required'
                })
            
            reader = PyPDF2.PdfReader(file)
            writer = PyPDF2.PdfWriter()
            total_pages = len(reader.pages)
            
            if total_pages == 0:
                return JsonResponse({
                    'success': False,
                    'error': 'The uploaded PDF has no pages'
                })
            
            # Parse page range
            pages_to_remove_list = []
            try:
                for part in pages_to_remove.split(','):
                    part = part.strip()
                    if '-' in part:
                        start, end = map(int, part.split('-'))
                        # Check if range is valid
                        if start < 1 or end > total_pages or start > end:
                            return JsonResponse({
                                'success': False,
                                'error': f'Invalid page range: {part}. Valid range is 1-{total_pages}'
                            })
                        pages_to_remove_list.extend(range(start, end + 1))
                    else:
                        page = int(part)
                        # Check if page is valid
                        if page < 1 or page > total_pages:
                            return JsonResponse({
                                'success': False,
                                'error': f'Invalid page number: {page}. Valid range is 1-{total_pages}'
                            })
                        pages_to_remove_list.append(page)
            except ValueError:
                return JsonResponse({
                    'success': False,
                    'error': 'Invalid page format. Use numbers and ranges (e.g., 1,3,5-7)'
                })
            
            # If we're trying to remove all pages, return an error
            if len(pages_to_remove_list) >= total_pages:
                return JsonResponse({
                    'success': False,
                    'error': 'Cannot remove all pages from the PDF'
                })
            
            # Add all pages except those to be removed
            for i, page in enumerate(reader.pages):
                if i + 1 not in pages_to_remove_list:
                    writer.add_page(page)
            
            output_path = os.path.join(settings.MEDIA_ROOT, 'modified.pdf')
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return JsonResponse({
                'success': True,
                'file_url': f'/media/modified.pdf'
            })
        except KeyError as e:
            return JsonResponse({
                'success': False,
                'error': f'Missing required field: {str(e)}'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({'error': 'Invalid request method'})

def register_view(request):
    if request.method == 'POST':
        form = CustomUserCreationForm(request.POST)
        print(f"Form data received: {request.POST}")
        print(f"Form is valid: {form.is_valid()}")
        if form.is_valid():
            try:
                user = form.save()
                login(request, user)
                messages.success(request, "Registration successful! You are now logged in.")
                return redirect('home')
            except Exception as e:
                print(f"Error saving user: {str(e)}")
                messages.error(request, f"Error during registration: {str(e)}")
        else:
            print(f"Form errors: {form.errors}")
            for field, errors in form.errors.items():
                for error in errors:
                    messages.error(request, f"{field}: {error}")
    else:
        form = CustomUserCreationForm()
    
    return render(request, 'pdf_tools/auth/register.html', {'form': form})

def login_view(request):
    if request.method == 'POST':
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            username = form.cleaned_data.get('username')
            password = form.cleaned_data.get('password')
            user = authenticate(username=username, password=password)
            if user is not None:
                login(request, user)
                messages.success(request, f"Welcome back, {username}!")
                next_url = request.GET.get('next', 'home')
                return redirect(next_url)
        messages.error(request, "Invalid username or password.")
    else:
        form = AuthenticationForm()
    
    return render(request, 'pdf_tools/auth/login.html', {'form': form})

def logout_view(request):
    logout(request)
    messages.success(request, "You have been logged out successfully.")
    return redirect('home')

@login_required
def my_documents(request):
    documents = ProcessedDocument.objects.filter(user=request.user).order_by('-created_at')
    context = {'documents': documents}
    return render(request, 'pdf_tools/user/my_documents.html', context)

@login_required
def document_history(request):
    history = DocumentHistory.objects.filter(user=request.user)
    context = {'history': history}
    return render(request, 'pdf_tools/user/document_history.html', context)

@login_required
def delete_document(request, document_id):
    document = get_object_or_404(ProcessedDocument, id=document_id, user=request.user)
    
    if request.method == 'POST':
        document.delete()
        messages.success(request, f"Document '{document.title}' deleted successfully.")
        return redirect('my_documents')
    
    return render(request, 'pdf_tools/user/confirm_delete.html', {'document': document})

@login_required
def toggle_favorite(request, document_id):
    if request.method == 'POST':
        document = get_object_or_404(ProcessedDocument, id=document_id, user=request.user)
        document.is_favorite = not document.is_favorite
        document.save()
        return JsonResponse({'status': 'success', 'is_favorite': document.is_favorite})
    
    return JsonResponse({'status': 'error', 'message': 'Invalid request'})

@login_required
def profile_view(request):
    user = request.user
    storage_used = sum(doc.file_size for doc in user.documents.all())
    storage_limit = user.profile.storage_limit
    storage_percent = (storage_used / storage_limit) * 100 if storage_limit > 0 else 0
    
    recent_documents = user.documents.order_by('-created_at')[:5]
    recent_history = user.history.order_by('-timestamp')[:5]
    
    context = {
        'user': user,
        'storage_used': storage_used,
        'storage_limit': storage_limit,
        'storage_percent': storage_percent,
        'recent_documents': recent_documents,
        'recent_history': recent_history
    }
    
    return render(request, 'pdf_tools/user/profile.html', context)

# Helper function to save processed files for authenticated users
def save_processed_file(user, file_path, title, document_type, operation, details=None):
    """
    Save a processed file to the user's document storage and create a history entry
    
    Args:
        user: User object
        file_path: Path to the processed file
        title: Document title
        document_type: Type of document (from ProcessedDocument.DOCUMENT_TYPES)
        operation: Operation performed (from DocumentHistory.OPERATIONS)
        details: Optional details about the operation
    
    Returns:
        The created ProcessedDocument instance
    """
    if not user.is_authenticated:
        return None
        
    # Get file size
    file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
    
    # Check if user has enough storage
    storage_used = sum(doc.file_size for doc in user.documents.all())
    if storage_used + file_size > user.profile.storage_limit:
        return None
    
    # Create document
    with open(file_path, 'rb') as f:
        document = ProcessedDocument(
            user=user,
            title=title,
            document_type=document_type,
            file_size=file_size
        )
        document.file.save(os.path.basename(file_path), f)
        document.save()
    
    # Create history entry
    history = DocumentHistory(
        user=user,
        operation=operation,
        output_document=document,
        details=details
    )
    history.save()
    
    return document
